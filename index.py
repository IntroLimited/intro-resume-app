import json
import os
import re
import subprocess
import tempfile
import urllib.request
import urllib.error
import urllib.parse
from http.server import BaseHTTPRequestHandler

NOTION_API_KEY = os.environ.get("NOTION_API_KEY", "")
NOTION_DATABASE_ID = os.environ.get("NOTION_DATABASE_ID", "036bdd7a61694c0e95450a26984e84c4")
TRELLO_DATABASE_ID = os.environ.get("TRELLO_DATABASE_ID", "b3e6ab7470994b3690c110b728f6593b")
GOOGLE_CLIENT_ID = os.environ.get("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.environ.get("GOOGLE_CLIENT_SECRET", "")
GOOGLE_REFRESH_TOKEN = os.environ.get("GOOGLE_REFRESH_TOKEN", "")

def notion_request(method, path, payload=None):
    url = f"https://api.notion.com/v1{path}"
    headers = {"Authorization": f"Bearer {NOTION_API_KEY}", "Content-Type": "application/json", "Notion-Version": "2022-06-28"}
    data = json.dumps(payload).encode() if payload else None
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())

def get_google_token():
    payload = {"client_id": GOOGLE_CLIENT_ID, "client_secret": GOOGLE_CLIENT_SECRET, "refresh_token": GOOGLE_REFRESH_TOKEN, "grant_type": "refresh_token"}
    req = urllib.request.Request("https://oauth2.googleapis.com/token", data=json.dumps(payload).encode(), headers={"Content-Type": "application/json"}, method="POST")
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())["access_token"]

def google_get(url, token):
    req = urllib.request.Request(url, headers={"Authorization": f"Bearer {token}"})
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())

def get_live_roles():
    roles = []
    cursor = None
    while True:
        payload = {"filter": {"property": "Status", "status": {"equals": "Live"}}, "page_size": 100}
        if cursor:
            payload["start_cursor"] = cursor
        data = notion_request("POST", f"/databases/{TRELLO_DATABASE_ID}/query", payload)
        for page in data.get("results", []):
            props = page.get("properties", {})
            name_prop = props.get("Name", {})
            role_name = "".join(t.get("plain_text", "") for t in name_prop.get("title", [])).strip()
            client_prop = props.get("Client", {})
            client_select = client_prop.get("select")
            client_name = client_select.get("name", "").strip() if client_select else ""
            if role_name and client_name:
                roles.append({"client": client_name, "role": role_name})
        if not data.get("has_more"):
            break
        cursor = data.get("next_cursor")
    roles.sort(key=lambda x: (x["client"].lower(), x["role"].lower()))
    return roles

def find_notion_candidate(candidate_name):
    try:
        data = notion_request("POST", f"/databases/{NOTION_DATABASE_ID}/query", {"filter": {"property": "title", "title": {"contains": candidate_name}}})
        results = data.get("results", [])
        target = candidate_name.lower().strip()
        for page in results:
            name_prop = (page.get("properties", {}).get("\ufeffName") or page.get("properties", {}).get("Name", {}))
            page_name = "".join(t.get("plain_text", "") for t in name_prop.get("title", [])).lower().strip()
            if page_name == target:
                return page["id"], page
        best, best_score = None, 0
        for page in results:
            name_prop = (page.get("properties", {}).get("\ufeffName") or page.get("properties", {}).get("Name", {}))
            page_name = "".join(t.get("plain_text", "") for t in name_prop.get("title", [])).lower().strip()
            score = sum(1 for w in target.split() if w in page_name.split())
            if score > best_score:
                best_score, best = score, page
        if best:
            return best["id"], best
    except Exception:
        pass
    return None, None

def get_candidate_data(page):
    props = page.get("properties", {})
    def rt(k): return "".join(t.get("plain_text","") for t in props.get(k,{}).get("rich_text",[]))
    def url(k): return props.get(k,{}).get("url","") or ""
    def status(k):
        s = props.get(k,{}).get("status"); return s.get("name","") if s else ""
    def ms(k): return ", ".join(i.get("name","") for i in props.get(k,{}).get("multi_select",[]))
    name_prop = (props.get("\ufeffName") or props.get("Name", {}))
    full_name = "".join(t.get("plain_text","") for t in name_prop.get("title",[])).strip()
    notes = rt("Notes")
    def section(notes, header):
        m = re.search(rf'{header}\s*\n(.*?)(?=\n[A-Z][A-Z ]+\n|$)', notes, re.DOTALL|re.IGNORECASE)
        return m.group(1).strip() if m else ""
    return {
        "full_name": full_name, "location": ms("Current Location"),
        "linkedin": url("LinkedIn"), "stage": status("Stage"),
        "basics": section(notes,"BASICS"), "strong_points": section(notes,"STRONG POINTS"),
        "potential_challenges": section(notes,"POTENTIAL CHALLENGES"), "compensation": section(notes,"COMPENSATION")
    }

def find_folder(token, parent_id, name):
    q = urllib.parse.quote(f"name='{name}' and mimeType='application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed=false")
    r = google_get(f"https://www.googleapis.com/drive/v3/files?q={q}&fields=files(id,name)", token)
    files = r.get("files",[])
    return files[0]["id"] if files else None

def find_presentation(token, folder_id):
    q = urllib.parse.quote(f"mimeType='application/vnd.google-apps.presentation' and '{folder_id}' in parents and trashed=false")
    r = google_get(f"https://www.googleapis.com/drive/v3/files?q={q}&fields=files(id,name,webViewLink)", token)
    files = r.get("files",[])
    return files[0] if files else None

def build_report(token, pres_id, candidate):
    # Download PPTX
    req = urllib.request.Request(f"https://www.googleapis.com/drive/v3/files/{pres_id}/export?mimeType=application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Authorization": f"Bearer {token}"})
    with urllib.request.urlopen(req) as resp:
        pptx_bytes = resp.read()

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        f.write(pptx_bytes)
        tmp_in = f.name
    tmp_out = tmp_in.replace('.pptx','_out.pptx')

    name_parts = candidate["full_name"].split()
    first = name_parts[0] if name_parts else candidate["full_name"]
    last = " ".join(name_parts[1:]) if len(name_parts) > 1 else ""

    script = f'''
import copy
from pptx import Presentation
prs = Presentation("{tmp_in}")
template_idx = 1 if len(prs.slides) > 1 else 0
template_slide = prs.slides[template_idx]
slide_layout = template_slide.slide_layout
new_slide = prs.slides.add_slide(slide_layout)
template_sp_tree = template_slide.shapes._spTree
new_sp_tree = new_slide.shapes._spTree
for i in range(len(new_sp_tree)-1,-1,-1):
    new_sp_tree.remove(new_sp_tree[i])
for child in template_sp_tree:
    new_sp_tree.append(copy.deepcopy(child))

def set_text(para, text):
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]: r.text = ""

data = {{
    "first": {repr(first)}, "last": {repr(last)},
    "location": {repr(candidate["location"])},
    "linkedin": {repr(candidate["linkedin"])},
    "stage": {repr(candidate["stage"])},
    "basics": {repr(candidate["basics"])},
    "strong_points": {repr(candidate["strong_points"])},
    "challenges": {repr(candidate["potential_challenges"])},
    "compensation": {repr(candidate["compensation"])},
}}

for shape in new_slide.shapes:
    if not shape.has_text_frame: continue
    full = shape.text_frame.text.upper()
    paras = shape.text_frame.paragraphs
    if ("CANDIDATE" in full and "NAME" in full) or full.strip() in ["NAME","CANDIDATE NAME"]:
        if len(paras) >= 2:
            set_text(paras[0], data["first"])
            set_text(paras[1], data["last"])
        elif paras: set_text(paras[0], data["first"]+" "+data["last"])
    elif "BASICS" in full:
        for i,p in enumerate(paras):
            t = p.text.upper().strip()
            if "BASICS" in t and i+1<len(paras): set_text(paras[i+1], data["basics"])
            elif "STRONG" in t and i+1<len(paras): set_text(paras[i+1], data["strong_points"])
            elif "POTENTIAL" in t and i+1<len(paras): set_text(paras[i+1], data["challenges"])
            elif "COMPENSATION" in t and i+1<len(paras): set_text(paras[i+1], data["compensation"])
            elif "STATUS" in t: set_text(p, "Status: "+data["stage"])
    elif "LOCATION" in full: set_text(paras[0], data["location"])
    elif "LINKEDIN" in full: set_text(paras[0], data["linkedin"] or "LinkedIn")

prs.save("{tmp_out}")
print("OK")
'''
    result = subprocess.run(['python3','-c',script], capture_output=True, text=True, timeout=60)
    if "OK" not in result.stdout:
        raise Exception(f"Slide error: {result.stderr[:200]}")

    with open(tmp_out,'rb') as f:
        data = f.read()
    req = urllib.request.Request(
        f"https://www.googleapis.com/upload/drive/v3/files/{pres_id}?uploadType=media",
        data=data,
        headers={"Authorization":f"Bearer {token}","Content-Type":"application/vnd.openxmlformats-officedocument.presentationml.presentation"},
        method="PATCH"
    )
    with urllib.request.urlopen(req) as resp:
        json.loads(resp.read())
    try: os.unlink(tmp_in); os.unlink(tmp_out)
    except: pass

class Handler(BaseHTTPRequestHandler):
    def do_GET(self):
        path = self.path.split('?')[0]
        if path == '/api/roles':
            try:
                roles = get_live_roles()
                self._json(200, {"roles": roles})
            except Exception as e:
                self._json(500, {"error": str(e)})
        elif path in ('/', '/index.html'):
            self._serve_file(os.path.join(os.path.dirname(__file__),'..','public','index.html'), 'text/html')
        else:
            self.send_error(404)

    def do_POST(self):
        if self.path.split('?')[0] == '/api/report':
            try:
                cl = int(self.headers.get('Content-Length',0))
                body = json.loads(self.rfile.read(cl))
                name = body.get('candidate_name','').strip()
                client = body.get('client','').strip()
                role = body.get('role','').strip()
                if not name: return self._json(400,{'error':'Candidate name required.'})
                if not client or not role: return self._json(400,{'error':'Client and role required.'})

                page_id, page = find_notion_candidate(name)
                if not page_id: return self._json(404,{'error':f'"{name}" not found in Notion.'})

                candidate = get_candidate_data(page)
                token = get_google_token()

                q = urllib.parse.quote(f"name='{client}' and mimeType='application/vnd.google-apps.folder' and trashed=false")
                r = google_get(f"https://www.googleapis.com/drive/v3/files?q={q}&fields=files(id,name)", token)
                client_folders = r.get("files",[])
                if not client_folders: return self._json(404,{'error':f'Client folder "{client}" not found in Drive.'})

                role_folder_id = find_folder(token, client_folders[0]["id"], role)
                if not role_folder_id: return self._json(404,{'error':f'Role folder "{role}" not found inside "{client}".'})

                pres = find_presentation(token, role_folder_id)
                if not pres: return self._json(404,{'error':f'No presentation found in "{client} → {role}".'})

                build_report(token, pres["id"], candidate)
                self._json(200,{'success':True,'deck_url':pres.get("webViewLink","")})

            except urllib.error.HTTPError as e:
                self._json(500,{'error':f'API error: {e.read().decode()[:200]}'})
            except Exception as e:
                self._json(500,{'error':str(e)})
        else:
            self.send_error(404)

    def _serve_file(self, path, ct):
        try:
            with open(path,'rb') as f: content = f.read()
            self.send_response(200)
            self.send_header('Content-Type', ct)
            self.send_header('Content-Length', str(len(content)))
            self.end_headers()
            self.wfile.write(content)
        except: self.send_error(404)

    def _json(self, status, data):
        body = json.dumps(data).encode()
        self.send_response(status)
        self.send_header('Content-Type','application/json')
        self.send_header('Content-Length',str(len(body)))
        self.send_header('Access-Control-Allow-Origin','*')
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, *args): pass
